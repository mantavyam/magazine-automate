# Install necessary libraries at the top
!pip install python-pptx
!pip install --upgrade python-docx

from google.colab import auth, drive
from pptx import Presentation
from docx import Document
from docx.shared import Pt
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
import os
import re

# Authenticate and mount Google Drive
auth.authenticate_user()
drive.mount('/content/drive')

# Function to check if text is in Hindi
def is_hindi(text):
    hindi_chars = [chr(i) for i in range(2304, 2432)]  # Unicode range for Hindi characters
    return any(char in hindi_chars for char in text)

# Function to extract text and table data from a PPTX file
def extract_text_and_tables_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    slides_data = []

    for slide in presentation.slides:
        slide_data = {"text": [], "tables": []}
        first_text = True

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        if first_text and not text.startswith("About") and 15 < len(text) < 120:
                            slide_data["text"].append(("headline", text))
                            first_text = False
                        else:
                            slide_data["text"].append(("paragraph", text))
            if shape.has_table:
                table_data = []
                table = shape.table
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                slide_data["tables"].append(table_data)

        slides_data.append(slide_data)

    return slides_data

# Function to create a DOCX file with extracted data
def create_docx_with_data(docx_path, slides_data):
    document = Document()

    for slide_index, slide in enumerate(slides_data):
        if slide_index > 0:
            slide_divider = document.add_paragraph("NEWSLIDE")  # Add slide divider
            slide_divider.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Add text data
        for text_type, text in slide["text"]:
            if text_type == "headline":
                document.add_heading(text, level=2)
            else:
                document.add_paragraph(text)

        # Add table data
        for table in slide["tables"]:
            doc_table = document.add_table(rows=len(table), cols=len(table[0]))
            doc_table.alignment = WD_TABLE_ALIGNMENT.CENTER

            for i, row in enumerate(table):
                for j, cell_text in enumerate(row):
                    cell = doc_table.cell(i, j)
                    cell.text = cell_text

            for row in doc_table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(10)

    document.save(docx_path)

# Function to format the DOCX file
def format_docx(docx_path):
    document = Document(docx_path)
    for paragraph in document.paragraphs:
        if is_hindi(paragraph.text) or 'kapil kathpal' in paragraph.text.lower():
            p = paragraph._element
            p.getparent().remove(p)
            p._p = p._element = None

    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                cell_paragraphs = list(cell.paragraphs)
                for paragraph in cell_paragraphs:
                    if is_hindi(paragraph.text) or 'kapil kathpal' in paragraph.text.lower():
                        p = paragraph._element
                        p.getparent().remove(p)
                        p._p = p._element = None

    document.save(docx_path)

# Function to identify categories and make them heading level 1
def identify_and_make_headings(docx_path):
    document = Document(docx_path)
    headings_style = document.styles['Heading1']
    paragraphs = list(document.paragraphs)

    for i, paragraph in enumerate(paragraphs):
        text = paragraph.text.strip()
        if len(text) < 15 and not text.startswith("About") and "NEWSLIDE" not in text:
            paragraph.style = headings_style

    document.save(docx_path)


# Main function
def main():
    # Get the PPTX file name from user input
    pptx_file_name = input("Enter the name of the PPTX file stored in Google Drive (e.g., '1-1-24.pptx'): ")
    pptx_file_path = os.path.join("/content/drive/My Drive/", pptx_file_name)

    # Extract text and table data from the PPTX file
    slides_data = extract_text_and_tables_from_pptx(pptx_file_path)

    # Create a new DOCX file and save it to Google Drive
    docx_file_name = f"{pptx_file_name.rsplit('.', 1)[0]}-RAW.docx"
    docx_file_path = os.path.join("/content/drive/My Drive/", docx_file_name)
    create_docx_with_data(docx_file_path, slides_data)

    # Format the newly created DOCX file
    format_docx(docx_file_path)

    # Identify and make categories heading level 1
    identify_and_make_headings(docx_file_path)

    print(f"Data extracted from '{pptx_file_name}' and saved into '{docx_file_name}' in Google Drive.")

# Execute main function
main()
